import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    NAVY = RGBColor(27, 42, 74)
    ORANGE = RGBColor(255, 107, 53)
    LIGHT_BG = RGBColor(248, 249, 250)
    DARK_TEXT = RGBColor(33, 37, 41)
    WHITE = RGBColor(255, 255, 255)
    GRAY_CARD = RGBColor(238, 242, 246)
    BORDER_GRAY = RGBColor(210, 215, 220)

    def set_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, title, module_name=""):
        # Header Bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        header.line.fill.background()

        # Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), Inches(13.333), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE
        line.line.fill.background()

        # Title Text
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12.0), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        if module_name:
            p0 = tf.paragraphs[0]
            p0.text = module_name.upper()
            p0.font.size = Pt(11)
            p0.font.bold = True
            p0.font.color.rgb = ORANGE
            p0.font.name = "Arial"

        p1 = tf.add_paragraph() if module_name else tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.font.name = "Arial"

    def add_card(slide, left, top, width, height, title="", content="", bg_color=GRAY_CARD, border_color=BORDER_GRAY):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)

        tb = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(width - 0.3), Inches(height - 0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        if title:
            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.size = Pt(15)
            p0.font.bold = True
            p0.font.color.rgb = NAVY
            p0.font.name = "Arial"
            p0.space_after = Pt(8)

        if content:
            lines = content.split('\n')
            for i, line in enumerate(lines):
                p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
                p.text = line
                p.font.size = Pt(12)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Arial"
                p.space_after = Pt(4)
        return shape

    def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
        num_rows = len(rows) + 1
        num_cols = len(headers)
        table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(left), Inches(top), Inches(width), Inches(height))
        table = table_shape.table

        if col_widths:
            for i, w in enumerate(col_widths):
                table.columns[i].width = Inches(w)

        # Style Headers
        for i, header_text in enumerate(headers):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            p = cell.text_frame.paragraphs[0]
            p.text = header_text
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.size = Pt(12)
            p.font.name = "Arial"
            p.alignment = PP_ALIGN.CENTER

        # Style Rows
        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else GRAY_CARD
                p = cell.text_frame.paragraphs[0]
                p.text = str(val)
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Arial"

    # ==================== SLIDE 1 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, NAVY)
    
    # Title Box
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "GIÁO ÁN TỔNG QUAN"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    p2.text = "CONTENT MARKETING TỔNG HỢP & BÀI BẢN BÁN HÀNG"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = "Arial"
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Mục tiêu: Nắm vững kiến thức marketing nền tảng, hành trình khách hàng, phễu bán hàng và tư duy tạo nội dung chuyển đổi cao."
    p3.font.size = Pt(15)
    p3.font.color.rgb = RGBColor(220, 225, 230)
    p3.font.name = "Arial"

    # ==================== SLIDE 2 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 2: Các Khái Niệm Cốt Lõi", "Module 1: Nền Tảng Marketing & Tư Duy Content")
    
    add_card(slide, 0.6, 1.5, 3.8, 5.3, "1. Marketing", "• Khái niệm: Quá trình nghiên cứu thị trường, phân tích đối thủ, xây dựng chiến lược, quảng bá, phân phối sản phẩm/dịch vụ và chăm sóc khách hàng.\n\n• Mục đích: Đáp ứng đúng nhu cầu, mong muốn của khách hàng để tạo ra giá trị & lợi nhuận bền vững.")
    add_card(slide, 4.76, 1.5, 3.8, 5.3, "2. Content Marketing", "• Khái niệm: Sáng tạo, xây dựng và phân phối nội dung có giá trị, phù hợp, nhất quán nhằm thu hút và giữ chân đúng đối tượng mục tiêu.\n\n• Mục đích: Nuôi dưỡng nhận thức, niềm tin, thiện cảm và thúc đẩy tăng trưởng doanh số.")
    add_card(slide, 8.93, 1.5, 3.8, 5.3, "3. Video Content", "• Khái niệm: Dạng content marketing dùng hình ảnh chuyển động để truyền tải thông điệp từ người bán đến khách hàng.\n\n• Sơ đồ dòng chảy truyền thông:\n[Thương hiệu (Brand)] → [Thông điệp (Content/Video)] → [Khách hàng (Customer)]")

    # ==================== SLIDE 3 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 3: Ba Tầng Giá Trị Sản Phẩm (Tư Duy 'Viết Cho Ai?')", "Module 1: Nền Tảng Marketing & Tư Duy Content")

    add_card(slide, 0.6, 1.5, 4.0, 5.3, "Tư Duy Định Hướng", "• Đừng hỏi: 'Viết sao cho hay?'\n• Hãy hỏi: 'Sản phẩm này đang giúp ai, trong chuyện gì?'\n\n• Nguyên lý: Một sản phẩm có nhiều kiểu người dùng với nhu cầu khác nhau. Viết chung chung sẽ không chạm vào ai.\n\n• Kết luận: Viết không rườm rà — chỉ cần đúng người.\n(Biết rõ người dùng → Thấy điểm khác biệt → Kể đúng thứ họ cần nghe)")

    headers = ["Tầng giá trị", "Bản chất", "Mục đích khi viết", "Vị trí"]
    rows = [
        ["1. Vật lý\n(Functional Value)", "Cung cấp rõ ràng, hữu hình, đo lường được: thành phần, công dụng, cách dùng.", "Để KH nhận ra: 'Ờ, cái này dùng được cho mình không?'", "Bảng thông tin sản phẩm"],
        ["2. Tinh thần\n(Emotional Value)", "Cảm xúc mang lại khi dùng: nhẹ bụng hơn, yên tâm hơn, đủ chất hơn...", "Để KH thấy: 'Hình như đúng cái mình đang cần.'", "Trải nghiệm thật"],
        ["3. Biểu tượng\n(Symbolic Value)", "Đại diện phong cách sống, hệ giá trị hướng tới (chọn vì giống kiểu người muốn trở thành).", "Để KH nhận ra: 'Ờ, cái này đúng gu mình.'", "Thói quen & lối sống"]
    ]
    add_table(slide, 4.8, 1.5, 7.9, 5.3, headers, rows, [1.5, 2.4, 2.5, 1.5])

    # ==================== SLIDE 4 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 4: Customer Journey (Hành Trình Khách Hàng)", "Module 1: Nền Tảng Marketing & Tư Duy Content")

    add_card(slide, 0.6, 1.4, 12.13, 1.0, "Khái niệm Customer Journey", "Toàn bộ trải nghiệm của khách hàng với thương hiệu từ lúc biết đến sản phẩm cho tới (và cả sau) khi mua hàng.")

    steps = [
        ("1. Nhận biết", "Awareness", "Biết đến thương hiệu lần đầu qua ads, social, bạn bè giới thiệu..."),
        ("2. Cân nhắc", "Consideration", "Tìm hiểu, so sánh, đọc review, xem giá, công dụng."),
        ("3. Mua hàng", "Purchase", "Quyết định mua hay không — giai đoạn chốt sale."),
        ("4. Trải nghiệm", "Experience", "Sử dụng sản phẩm, đánh giá chất lượng, cảm xúc sau mua."),
        ("5. Trung thành", "Loyalty", "Quay lại mua tiếp, giới thiệu bạn bè nếu hài lòng.")
    ]

    for idx, (title_vn, title_en, desc) in enumerate(steps):
        left = 0.6 + idx * 2.48
        add_card(slide, left, 2.7, 2.3, 4.1, f"{title_vn}\n({title_en})", desc)

    # ==================== SLIDE 5 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 5: Sales Funnel & Mô Hình TOFU – MOFU – BOFU", "Module 1: Nền Tảng Marketing & Tư Duy Content")

    add_card(slide, 0.6, 1.4, 12.13, 1.2, "Bản Chất Sales Funnel & Vai Trò Của Content", "• Quá trình lọc/chuyển đổi: Biết đến lần đầu → Tò mò/tìm hiểu → Tương tác sâu → Xem xét/quyết định → Mua hàng.\n• Vai trò Content: Khách hàng nằm ở nhiều giai đoạn khác nhau; Content Marketing tạo nội dung đúng — đúng người — đúng giai đoạn để dẫn dắt chuyển đổi.")

    headers = ["Tầng phễu", "Mục tiêu", "Trạng thái KH", "Câu hỏi cần trả lời", "Gợi ý Content"]
    rows = [
        ["TOFU\n(Top of Funnel)", "Thu hút chú ý, tạo nhận biết thương hiệu", "Đang có vấn đề, tìm giải pháp", "Vấn đề/nhu cầu là gì? Vì sao xuất hiện?", "Blog, social post, infographic, ebook, video, checklist, newsletter..."],
        ["MOFU\n(Middle of Funnel)", "Thuyết phục giải pháp phù hợp", "Đang cân nhắc nhiều lựa chọn", "Tiêu chí so sánh? Có giải pháp nào tốt hơn?", "So sánh khách quan, feedback thật, hé lộ quy trình sản xuất..."],
        ["BOFU\n(Bottom of Funnel)", "Thúc đẩy ra quyết định mua", "Đã sẵn sàng, cần cú hích cuối", "Lý do quyết định mua? Tiêu chí quan trọng nhất?", "Khuyến mãi + CTA mạnh, chứng thực từ khách hàng thật (testimonial)..."]
    ]
    add_table(slide, 0.6, 2.8, 12.13, 4.0, headers, rows, [2.0, 2.3, 2.3, 2.5, 3.03])

    # ==================== SLIDE 6 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 6: Phân Biệt Vấn Đề — Pain Point — Customer Insight", "Module 2: Customer Insight & Kỹ Năng Viết Content")

    add_card(slide, 0.6, 1.4, 12.13, 1.0, "Định nghĩa Customer Insight", "Là 'sự thật ngầm hiểu' xuất phát từ niềm tin và hành vi, thúc đẩy quá trình suy nghĩ, quyết định và hành động mua hàng của khách hàng.")

    headers = ["Khái niệm", "Là gì?", "Biểu hiện", "Ghi chú"]
    rows = [
        ["Vấn đề", "Những gì KH gặp trong thực tế, nhận thức được và nói ra được.", "Nói ra được rõ ràng.", "Ở mức bề mặt, dễ khai thác nhất."],
        ["Pain Point", "Nỗi đau sâu hơn (cảm xúc tiêu cực, nỗi lo, sự khó chịu, nỗi sợ).", "Cảm giác lo lắng, áy náy, khó chịu.", "Sâu hơn vấn đề một bậc."],
        ["Customer Insight", "Sự thật sâu bên trong hành vi/cảm xúc; nếu 'chạm đúng' sẽ tạo đồng cảm và chuyển đổi mạnh.", "Ít khi nói ra, nhưng ai nói hộ thì gật gù.", "Không hiển lộ rõ, cần nghiên cứu kỹ."]
    ]
    add_table(slide, 0.6, 2.6, 12.13, 4.2, headers, rows, [2.0, 3.8, 3.3, 3.03])

    # ==================== SLIDE 7 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 7: Mô Hình Viết Content AIDA", "Module 3: Công Thức Viết Content & Kỹ Thuật Hook")

    cards_aida = [
        ("A — Attention (Gây chú ý)", "• Bước quan trọng nhất quyết định người đọc dừng lại.\n• Tiêu đề là ưu tiên số 1.\n• Dùng con số, câu hỏi, sự thật gây sốc, từ ngữ giật tít."),
        ("I — Interest (Tạo thích thú)", "• Cung cấp thông tin giúp KH giải đáp vấn đề/thắc mắc chưa có lời giải.\n• Nội dung dạng hướng dẫn, giải đáp thắc mắc cụ thể."),
        ("D — Desire (Thúc đẩy khao khát)", "• Thể hiện ưu điểm, lợi thế cạnh tranh, lợi ích cụ thể.\n• Khai thác triệt để insight để chạm đúng nhu cầu KH."),
        ("A — Action (Kêu gọi hành động)", "• Chốt sale với CTA văn phong mạnh mẽ, thuyết phục, hình thức rõ ràng.\n• Ví dụ: 'Mua ngay', 'Nhận ưu đãi hôm nay', 'Nhắn tin ngay'...")
    ]

    for idx, (title, content) in enumerate(cards_aida):
        row = idx // 2
        col = idx % 2
        left = 0.6 + col * 6.2
        top = 1.5 + row * 2.7
        add_card(slide, left, top, 5.9, 2.5, title, content)

    # ==================== SLIDE 8 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 8: Mô Hình Viết Content PAS & FAB", "Module 3: Công Thức Viết Content & Kỹ Thuật Hook")

    add_card(slide, 0.6, 1.5, 5.9, 5.3, "Mô Hình PAS (Dành Cho KH Đã Nhận Thức)", "• P — Problem (Vấn đề): Xác định đúng nỗi đau/tình trạng KH gặp phải (mở đầu bằng tin xã hội chấn động hoặc câu chuyện nhức nhối).\n\n• A — Agitate (Khuấy động): Làm trầm trọng nỗi đau lên nhiều lần để KH thấy cần xử lý ngay.\n\n• S — Solve (Giải pháp): Khẳng định sản phẩm là thứ duy nhất hóa giải được vấn đề.\n\n👉 Tâm lý học: Đáp ứng câu hỏi 'Tại sao tôi nên dùng sản phẩm này?'")
    add_card(slide, 6.8, 1.5, 5.9, 5.3, "Mô Hình FAB (Trọng Tâm Ở Benefits)", "• F — Features (Tính năng): Thông số, thành phần, công dụng sản phẩm có thể làm.\n\n• A — Advantages (Ưu điểm): Điểm vượt trội, lợi thế cạnh tranh so với đối thủ.\n\n• B — Benefits (Lợi ích): Lợi ích thực tế KH nhận được khi sử dụng.\n\n👉 Trọng tâm: Tâm lý KH ra quyết định nhanh hơn khi lợi ích cá nhân được đề cao rõ ràng.")

    # ==================== SLIDE 9 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 9: Kỹ Thuật Viết Tiêu Đề Thu Hút (Hook)", "Module 3: Công Thức Viết Content & Kỹ Thuật Hook")

    add_card(slide, 0.6, 1.5, 5.0, 5.3, "Tầm Quan Trọng & Nguyên Tắc", "• Tầm quan trọng:\nTiêu đề hấp dẫn → Tăng CTR → Tăng Lead → Tăng chuyển đổi.\n\n• 4 Nguyên tắc cốt lõi:\n1. Rõ ràng & súc tích.\n2. Sử dụng từ khóa hiệu quả (khớp search intent).\n3. Khơi gợi tò mò & cảm xúc (công thức 'hoài nghi – bí mật – sự thật').\n4. Độc đáo, mới mẻ & Phù hợp đối tượng mục tiêu.")
    add_card(slide, 5.9, 1.5, 6.8, 5.3, "9 Kỹ Thuật Viết Tiêu Đề Cụ Thể", "1. Sử dụng con số để tăng tính cụ thể, hấp dẫn.\n2. Đặt câu hỏi kích thích tò mò.\n3. Dùng từ ngữ mạnh mẽ, mang tính hành động.\n4. Tạo sự đối lập hoặc nghịch lý.\n5. Kể một câu chuyện ngắn.\n6. Trích dẫn câu nói của người nổi tiếng.\n7. Tạo cảm giác khan hiếm hoặc cấp bách.\n8. Dùng lời hứa để tạo niềm tin.\n9. Áp dụng các hot trend đang được quan tâm.")

    # ==================== SLIDE 10 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 10: Bài Tập Thực Hành (Buổi 3)", "Module 3: Công Thức Viết Content & Kỹ Thuật Hook")

    add_card(slide, 0.6, 1.6, 12.13, 2.4, "Nhiệm vụ 1: Viết bài Content Hoàn Chỉnh", "• Đề bài: Viết 01 bài content quảng cáo cho sản phẩm Bánh rau củ Enzy hoặc Gia vị rắc cơm.\n• Yêu cầu: Áp dụng chuẩn cấu trúc theo 2 mô hình đã học: AIDA và PAS.\n• Tiêu chí chấm: Xác định rõ Insight, Hook ấn tượng, Logic diễn giải chặt chẽ, CTA rõ ràng.")
    add_card(slide, 0.6, 4.3, 12.13, 2.4, "Nhiệm vụ 2: Sáng Tạo Tiêu Đề Thu Hút", "• Đề bài: Viết 05 tiêu đề hấp dẫn cho các sản phẩm: Hạt nêm chay, Rắc cơm, Bánh rau củ.\n• Yêu cầu: Áp dụng linh hoạt các kỹ thuật (Con số, Câu hỏi, Khan hiếm, Trend, Lời hứa...).\n• Thách thức: Đảm bảo vừa khơi gợi tò mò vừa giữ tính chân thật của sản phẩm.")

    # ==================== SLIDE 11 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 11: Quy Trình Sản Xuất Video Content (Bước 1 - 3)", "Module 4: Sản Xuất Video Content Chuyển Đổi")

    add_card(slide, 0.6, 1.5, 3.8, 5.3, "Bước 1: Xác Định Khách Hàng", "Ví dụ: Thực phẩm chay Enzy\n• Nhân khẩu học: Nữ (chủ yếu), 25–45t; Thu nhập trung bình – khá; Sống tại TP lớn/toàn quốc; Đã có gia đình/chăm sóc người thân.\n• Sở thích: Nấu ăn món chay, xem video mẹo bếp, tìm nguyên liệu sạch, tham gia cộng đồng sức khỏe.")
    add_card(slide, 4.76, 1.5, 3.8, 5.3, "Bước 2: Xác Định Pain Point & Insight", "• Pain Point: Nỗi lo thực phẩm bẩn, tốn thời gian chế biến món chay, nêm nếm không chuẩn vị.\n• Insight: Muốn chăm sóc gia đình bằng bữa ăn lành mạnh nhưng cần giải pháp tiện lợi, chuẩn vị nhanh chóng.\n• Lập bảng đối chiếu Pain Point ↔ Insight.")
    add_card(slide, 8.93, 1.5, 3.8, 5.3, "Bước 3: Phân Tích Hành Trình & Phễu", "Lập bảng Ma trận Content:\n• Giai đoạn Nhận thức: Vấn đề & lý do xuất hiện.\n• Giai đoạn Đánh giá/So sánh: Tiêu chí so sánh giải pháp.\n• Giai đoạn Quyết định: Tiêu chí cốt lõi để mua hàng.\n→ Insight → Hành vi → Gợi ý Content tương ứng.")

    # ==================== SLIDE 12 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 12: Quy Trình Video (Bước 4 - 6) & Form Kịch Bản", "Module 4: Sản Xuất Video Content Chuyển Đổi")

    add_card(slide, 0.6, 1.5, 5.0, 5.3, "Các Bước Tiếp Theo (Bước 4 - 6)", "• Bước 4: Chọn dạng content phù hợp (Review, Unboxing, Story, Chuyên gia...).\n\n• Bước 5: Viết kịch bản chi tiết:\n⚡ 3 giây đầu quyết định video 'win' hay không (Hook phải đủ mạnh).\n• CTR tăng → Lượt truy cập/inbox tăng → Tăng doanh thu.\n• CTR tăng → Thuật toán ưu tiên hiển thị → CPC giảm → Tiết kiệm ngân sách ads.\n\n• Bước 6: Chỉnh sửa & Tối ưu video.")

    headers = ["Đặt tên Video", "Kịch bản (Thoại/Hành động)", "Text Video", "Caption", "CTA / Tiêu đề"]
    rows = [
        ["Hook 3s", "Tình huống giật gân / Câu hỏi nhức nhối", "TEXT NỔI BẬT", "Mở đầu giật tít", "Dừng lướt ngay!"],
        ["Thân bài", "Diễn giải nỗi đau & Giới thiệu giải pháp", "Tính năng / Lợi ích chính", "Thông tin chi tiết", "Xem ngay bên dưới"],
        ["Kết bài", "Đưa ra ưu đãi & Hướng dẫn hành động", "ƯU ĐÃI HOM NAY", "Link / Hotline", "NHẮN TIN MUA NGAY"]
    ]
    add_table(slide, 5.9, 1.5, 6.8, 5.3, headers, rows, [1.3, 1.8, 1.3, 1.2, 1.2])

    # ==================== SLIDE 13 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 13: Bảy Dạng Video Content (Dạng 1 đến 4)", "Module 4: Sản Xuất Video Content Chuyển Đổi")

    v_types = [
        ("1. Dạng Truyền Thống", "Vấn đề (Hook) → Giới thiệu SP (công dụng, thành phần, lợi ích, công nghệ) → Điểm độc đáo (USP) → CTA."),
        ("2. Dạng Review", "Người thật/creator giới thiệu → Quá trình trải nghiệm trước – sau → Cảm nhận thực tế → So sánh nhẹ → CTA."),
        ("3. Dạng Unboxing", "Hook mở đầu mạnh mẽ → Quá trình mở hộp sản phẩm → Cảm nhận đầu tiên về hình thức & chất lượng → CTA."),
        ("4. Dạng Story", "Dẫn dắt bằng câu hook drama/cảm xúc → Tình huống chuyển biến → Kết thúc chạm cảm xúc → Lồng ghép SP/CTA.")
    ]

    for idx, (title, content) in enumerate(v_types):
        row = idx // 2
        col = idx % 2
        left = 0.6 + col * 6.2
        top = 1.5 + row * 2.7
        add_card(slide, left, top, 5.9, 2.5, title, content)

    # ==================== SLIDE 14 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 14: Bảy Dạng Video (5-7) & Ứng Dụng AI", "Module 4: Sản Xuất Video Content Chuyển Đổi")

    add_card(slide, 0.6, 1.5, 3.8, 3.2, "5. Dạng Chuyên Gia", "Chuyên gia/bác sĩ/dinh dưỡng xuất hiện → Giải thích vấn đề → Đưa giải pháp an toàn → Đề cập SP (CTA).")
    add_card(slide, 4.76, 1.5, 3.8, 3.2, "6. Testimonial / Feedback", "Cắt feedback thật từ khách hàng → Diễn giải chi tiết lời chứng thực → Thúc đẩy niềm tin → CTA.")
    add_card(slide, 8.93, 1.5, 3.8, 3.2, "7. Dạng So Sánh", "Hook về sự khác nhau → Nêu quan điểm phân tích → Nhấn mạnh USP sản phẩm → CTA.")

    add_card(slide, 0.6, 4.9, 12.13, 1.9, "🤖 Ứng Dụng AI Trong Sản Xuất Content", "• Tikviral: Hỗ trợ khai thác dữ liệu, nghiên cứu trend TikTok nhanh chóng.\n• Gemini, ChatGPT: Hỗ trợ lên ý tưởng, viết kịch bản chi tiết, tạo nội dung đa kênh.\n• Kalodata: Hỗ trợ phân tích dữ liệu bán hàng & thị trường thương mại điện tử.")

    # ==================== SLIDE 15 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 15: Quy Trình Content Marketing 6 Bước Khép Kín", "Module 5: Quy Trình & Tư Duy Làm Content Marketing")

    steps_6 = [
        ("1. Research", "Giai đoạn 1", "Nghiên cứu khách hàng & đối thủ sâu sắc."),
        ("2. Ideation", "Giai đoạn 1", "Lên ý tưởng & Lập kế hoạch nội dung."),
        ("3. Creation", "Giai đoạn 1", "Sáng tạo & Sản xuất nội dung đa định dạng."),
        ("4. Promotion", "Giai đoạn 2", "Phân phối nội dung trên các kênh mục tiêu."),
        ("5. Convert Lead", "Giai đoạn 2", "Tối ưu hóa tỷ lệ chuyển đổi thành Lead."),
        ("6. Measurement", "Giai đoạn 2", "Đo lường hiệu quả & Tối ưu vòng tiếp theo.")
    ]

    for idx, (title, g, desc) in enumerate(steps_6):
        left = 0.6 + idx * 2.05
        add_card(slide, left, 1.6, 1.9, 4.0, title, f"[{g}]\n\n{desc}")

    add_card(slide, 0.6, 5.8, 12.13, 1.0, "🔄 Nguyên Lý Khép Kín", "Kết quả đo lường ở Bước 6 sẽ làm dữ liệu đầu vào cho Bước 1 để liên tục tối ưu hiệu quả chiến dịch.")

    # ==================== SLIDE 16 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 16: Tổng Quan Phễu Chuyển Đổi & Định Dạng Content", "Module 5: Quy Trình & Tư Duy Làm Content Marketing")

    add_card(slide, 0.6, 1.4, 12.13, 0.9, "Vai Trò Của Content & Kênh Phân Phối", "• Content: Giúp đưa khách hàng 'xuống' các tầng phễu dưới bằng cách thay đổi thái độ & tâm lý.\n• Kênh phân phối: Đưa đúng content tới đúng người, đúng thời điểm.")

    headers = ["Tầng phễu", "Vai trò chính", "Định dạng nội dung phù hợp"]
    rows = [
        ["TOFU — Awareness\n(Nhận thức)", "Tăng độ nhận diện thương hiệu sâu rộng.", "Blog post, social post, infographic, digital magazine, ebook, video, report, checklist, newsletter..."],
        ["MOFU — Consideration\n(Cân nhắc)", "Cung cấp thông tin chuyên sâu, tăng uy tín, nuôi dưỡng lòng tin.", "Blog post, email, webinar, video demo, tài nguyên hữu ích, quiz, khảo sát..."],
        ["BOFU — Conversion\n(Chuyển đổi)", "Thúc đẩy & khuyến khích hành vi mua hàng ngay.", "Demo, customer story, spec sheet, webinar, testimonial, review, landing page, case study..."]
    ]
    add_table(slide, 0.6, 2.5, 12.13, 4.3, headers, rows, [2.5, 3.5, 6.13])

    # ==================== SLIDE 17 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 17: Xu Hướng Content 2025–2026", "Module 6: Xu Hướng & Tư Duy Đọc-Xem Của Khách Hàng")

    add_card(slide, 0.6, 1.5, 3.8, 5.3, "1. Short-form Video Chân Thật", "• Quảng cáo, review độc thoại, nhảy trend đang bão hòa.\n\n• Lên ngôi bền vững: Content dạng nhật ký, chia sẻ trải nghiệm thật, behind-the-scenes, tiểu phẩm đời thường & content giáo dục.")
    add_card(slide, 4.76, 1.5, 3.8, 5.3, "2. Content Tâm Sự Trên Threads", "• Nơi lưu lại câu chuyện, trải nghiệm cá nhân.\n\n• Phổ biến: Status ngắn mang tính self-reflection, self-healing, góc nhìn cá nhân sâu sắc thu hút Gen Z.")
    add_card(slide, 8.93, 1.5, 3.8, 5.3, "3. Experience Sharing (Bite-size)", "• Chia sẻ kinh nghiệm làm việc thực tế.\n\n• Định dạng: Dạng ngắn (bite-size), dễ đọc, dễ tiêu thụ, gần gũi và mang lại giá trị thực tế áp dụng ngay.")

    # ==================== SLIDE 18 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 18: Khái Niệm Readability & Watchability", "Module 6: Xu Hướng & Tư Duy Đọc-Xem Của Khách Hàng")

    add_card(slide, 0.6, 1.5, 5.9, 2.6, "Readability (Khả năng đọc — Thang 0-10)", "• Định nghĩa: Mức độ khách hàng sẵn sàng đọc & nghiền ngẫm nội dung văn bản.\n• Ngành Readability cao: Cần thời gian nghiên cứu, thông tin lớn (Tài chính, Giáo dục, BĐS...).")
    add_card(slide, 6.8, 1.5, 5.9, 2.6, "Watchability (Khả năng xem — Thang 0-10)", "• Định nghĩa: Mức độ nội dung cần trực quan hóa qua video, hình ảnh, âm thanh.\n• Ngành Watchability cao: Cần hướng dẫn tutorial, minh họa thị giác (Thời trang, Ấm thực, Du lịch...).")
    add_card(slide, 0.6, 4.3, 12.13, 2.5, "Phân Tích Ngành & Tư Duy Chiến Lược", "• Ví dụ ngành Thời trang: Readability thấp (1–3 điểm); KH ít 'đọc' mà chủ yếu xem ảnh, search Google, xem influencer.\n• Kết luận: Xác định đúng Readability/Watchability giúp quyết định viết gì, cho ai, nông hay sâu, dài hay ngắn và phân phối kênh nào.")

    # ==================== SLIDE 19 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 19: Case Study Phễu Theo Đặc Thù Ngành Hàng", "Module 6: Xu Hướng & Tư Duy Đọc-Xem Của Khách Hàng")

    add_card(slide, 0.6, 1.5, 5.9, 5.3, "Case 1: Ngành Dầu Gội (FMCG)", "Watchability Trung bình-Thấp | Readability Cao\n\n• Đặc điểm KH: Không thích đọc quá nhiều về SP, bị thu hút bởi ảnh đẹp, người nổi tiếng, thành phần minh họa.\n• Triển khai Phễu:\n- TOFU: Làm rất mạnh (Awareness) để KH nhớ thương hiệu.\n- MOFU: Bỏ qua hoặc làm rất ít.\n- BOFU: Kích thích mua bằng khuyến mãi, minigame, giveaway.")
    add_card(slide, 6.8, 1.5, 5.9, 5.3, "Case 2: Ngành Giáo Dục (Brand Course)", "Watchability Trung bình | Readability Cao\n\n• Đặc điểm KH: Sẵn sàng đọc sâu để nghiền ngẫm kiến thức trước khi mua.\n• Triển khai Phễu:\n- TOFU: Bài viết chia sẻ kiến thức chuyên ngành sâu.\n- MOFU: Livestream, Event giải đáp thắc mắc chuyên sâu.\n- BOFU: Landing page chi tiết khóa học, ads, product demo.")

    # ==================== SLIDE 20 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 20: Mô Hình 3C Trong Content Strategy", "Module 7: Mô Hình 3C, Strategy Canvas & Value Matrix")

    add_card(slide, 0.6, 1.5, 5.0, 5.3, "3 Vòng Tròn Giao Nhau", "• Bài toán cốt lõi: KH sẵn sàng chi trả bao nhiêu cho vấn đề? Mình giải quyết tốt tới đâu?\n\n1. Consumer's problem: Vấn đề/nỗi đau của khách hàng.\n2. Company's solution: Giải pháp của sản phẩm công ty.\n3. Competitor's solution: Giải pháp của đối thủ cạnh tranh.")
    add_card(slide, 5.9, 1.5, 6.8, 5.3, "Phân Tích POD & POP", "• Point of Difference (POD - Điểm khác biệt):\nNơi giải pháp của bạn vượt trội đối thủ → BẮT BUỘC PHẢI THẮNG.\n\n• Point of Parity (POP - Điểm tương đồng):\nNơi phải có mặt để không bị bỏ lại trên thị trường.\n\n⚠️ Lưu ý: Vấn đề phải cấp bách & Giải pháp phải có tính vượt trội rõ rệt. Nếu thiếu cả 2, Content sẽ rất khó bán hàng.")

    # ==================== SLIDE 21 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 21: Mô Hình Strategy Canvas", "Module 7: Mô Hình 3C, Strategy Canvas & Value Matrix")

    add_card(slide, 0.6, 1.5, 12.13, 2.0, "Định Nghĩa Strategy Canvas", "• Công cụ liệt kê các tiêu chí cạnh tranh trong ngành (ví dụ: ease of use, check-in, delays, safety, price...).\n• Vẽ đường biểu diễn mức độ đáp ứng của thương hiệu mình so với đối thủ để tìm điểm nổi trội.")
    add_card(slide, 0.6, 3.7, 5.9, 3.1, "Trường Hợp 1: Có Điểm Nổi Trội", "• Tập trung toàn bộ Key Message và các tuyến Content xoay quanh điểm nổi trội đó.\n• Khai thác triệt để trên mọi kênh truyền thông để khắc họa dấu ấn.")
    add_card(slide, 6.8, 3.7, 5.9, 3.1, "Trường Hợp 2: Không Có Điểm Nổi Trội", "• Bắt buộc phải cạnh tranh bằng giá (Price War).\n• HOẶC phải tập trung xây dựng năng lực cốt lõi / lợi thế cạnh tranh hoàn toàn mới (Root Strength).")

    # ==================== SLIDE 22 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 22: Ma Trận Giá Trị (Value Matrix)", "Module 7: Mô Hình 3C, Strategy Canvas & Value Matrix")

    add_card(slide, 0.6, 1.3, 12.13, 0.8, "Công Thức Cốt Lõi", "Nhóm Khách Hàng Tốt Nhất + Pain Point + Giá Trị Sản Phẩm = Key Message (Thông Điệp Truyền Thông)")

    headers = ["Nhóm KH", "Pain Point (Nỗi đau)", "Giá trị sản phẩm", "Thông điệp truyền thông (Key Message)"]
    rows = [
        ["Người trẻ đi làm\n(25–35t)", "Lương ổn định nhưng dễ tiêu hết; sợ rủi ro, sợ bị lừa.", "Đầu tư & BH đơn giản, tích lũy từ vốn nhỏ.", "'Tự do tài chính không đến từ lương cao — mà từ sự chuẩn bị thông minh sớm.'"],
        ["Gia đình trẻ\nco con nhỏ", "Chi phí tăng nhanh; sợ rủi ro trụ cột; muốn lo cho con.", "Hỗ trợ tài chính bất trắc; quỹ tích lũy học vấn.", "'An tâm cho con — Đầu tư linh hoạt, bảo vệ toàn diện.'"],
        ["Người trung niên\n(45–60t)", "Sợ gánh nặng con cái; sợ tốn tiền khi bệnh nặng.", "Bảo vệ toàn diện bệnh nặng; đầu tư an toàn.", "'Bảo vệ thu nhập — Đầu tư an toàn cho tuổi nghỉ hưu.'"],
        ["Lao động tự do", "Không BHXH; mất thu nhập khi bệnh; phí không cố định.", "Đóng linh hoạt, rút online 24/7, bảo vệ dài hạn.", "'Làm chủ cuộc sống tự do — vẫn có điểm tựa tài chính khi cần.'"]
    ]
    add_table(slide, 0.6, 2.2, 12.13, 4.7, headers, rows, [2.0, 3.0, 3.0, 4.13])

    # ==================== SLIDE 23 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 23: Mirage Content & Nguyên Nhân Không Ra Lead", "Module 8: Content Mapping & Chuyển Đổi Lead")

    add_card(slide, 0.6, 1.5, 5.9, 5.3, "Hiện Tượng Mirage Content (Nội Dung Ảo)", "• Trông có vẻ ổn, kéo traffic tốt nhưng KHÔNG CÓ LEAD vì thiếu tính cụ thể.\n\n• Ví dụ 1: 'Làm sao để startup thu hút khách?' → Quá rộng, thuộc phạm vi một cuốn sách.\n• Ví dụ 2: 'Phương pháp học tiếng Anh' → Quá chung chung, không nổi bật giữa hàng ngàn bài trên mạng.\n\n👉 Kết luận: Content chung chung chỉ tăng awareness chứ không chuyển đổi.")
    add_card(slide, 6.8, 1.5, 5.9, 5.3, "Customer-Centric Content Marketing", "• Specific Customer: Biết rõ khách hàng mục tiêu cụ thể là ai.\n• Specific Pain: Viết cho nỗi đau cấp bách KH sẵn sàng chi trả để giải quyết.\n\n💡 CÔNG THỨC: Customer Pain + USP (cụ thể) = Content ra Lead\n\n• Nguyên nhân khác: Content không khớp phễu, không nổi bật USP, giật tít rỗng, CTA sai/mơ hồ.")

    # ==================== SLIDE 24 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 24: Content Mapping & Bản Bồi Phễu", "Module 8: Content Mapping & Chuyển Đổi Lead")

    add_card(slide, 0.6, 1.5, 5.0, 5.3, "Định Nghĩa & Lợi Ích", "• Định nghĩa: Hoạt động lên kế hoạch nội dung nhằm đưa đúng nội dung đến đúng người vào đúng thời điểm.\n\n• Lợi ích:\n1. What to do: Biết KH cần gì ở từng giai đoạn & cách tiếp cận.\n2. What missing: Biết nội dung/kênh phân phối nào còn thiếu để bù đắp.")

    steps_cm = [
        ("Bước 1", "Xác định khách hàng mục tiêu cụ thể."),
        ("Bước 2", "Xác định mục tiêu của từng giai đoạn (Stage)."),
        ("Bước 3", "Phân tích đặc trưng khách hàng từng giai đoạn."),
        ("Bước 4", "Phân tích đặc trưng nội dung từng giai đoạn."),
        ("Bước 5", "Mapping nội dung vào phễu content bán hàng.")
    ]
    for idx, (b, d) in enumerate(steps_cm):
        top = 1.5 + idx * 1.05
        add_card(slide, 5.9, top, 6.8, 0.95, b, d)

    # ==================== SLIDE 25 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 25: Đặc Điểm Khách Hàng Ở Hai Đầu Phễu", "Module 8: Content Mapping & Chuyển Đổi Lead")

    headers = ["Đặc điểm", "TOFU (Tầng trên cùng)", "BOFU (Tầng dưới cùng)"]
    rows = [
        ["Trạng thái KH", "Chưa sẵn sàng chi trả, tiềm năng về nhân khẩu học.", "Sẵn sàng chi trả, nhu cầu cấp bách."],
        ["Nhiệm vụ Content", "Tăng nhận diện thương hiệu, tạo tương tác rộng.", "Giúp ra quyết định mua hàng nhanh nhất."],
        ["Cảnh báo ⚠️", "Tránh đăng nội dung kêu gọi mua hàng quá sớm.", "Tránh đăng tin tức chung chung gây xao nhãng."]
    ]
    add_table(slide, 0.6, 1.5, 12.13, 3.2, headers, rows, [2.5, 4.8, 4.83])

    add_card(slide, 0.6, 4.9, 12.13, 1.9, "💡 Chiến Lược Cân Bằng TOFU & BOFU", "• Nếu chỉ làm TOFU: Traffic rất lớn nhưng KHÔNG CHUYỂN ĐỔI ĐƯỢC AI.\n• Nếu chỉ làm BOFU: Tương tác thấp, tệp công chúng quá nhỏ không phát triển được.\n→ Cần kết hợp hài hòa cả 2 tầng trong kế hoạch Content Mapping.")

    # ==================== SLIDE 26 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 26: Chiến Lược Phân Phối Trên Facebook", "Module 9: Phân Phối Nội Dung Theo Nền Tảng")

    add_card(slide, 0.6, 1.5, 3.8, 5.3, "1. Timing & Frequency", "• Đăng giờ tối ưu dựa trên data hành vi (Meta Business Suite).\n\n• Tần suất: 1–3 bài/ngày để tăng nhận diện nhưng không gây phiền.")
    add_card(slide, 4.76, 1.5, 3.8, 5.3, "2. Format Optimization", "• Short-form (reels, ảnh đơn) cho Awareness.\n• Long-form (album, carousel) cho Consideration.\n• Ưu tiên Video, caption thu hút, tối ưu mobile, thumbnail và CTA rõ.")
    add_card(slide, 8.93, 1.5, 3.8, 5.3, "3. Strategic Distribution", "• Seeding: Tập trung dồn lực comment/share trong 2 giờ đầu.\n• Chia sẻ chéo: Group, fanpage phụ, trang cá nhân.\n• Community: Ghim comment thú vị, phản hồi duy trì hội thoại.")

    # ==================== SLIDE 27 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 27: Thuật Toán Phân Phối TikTok", "Module 9: Phân Phối Nội Dung Theo Nền Tảng")

    add_card(slide, 0.6, 1.5, 5.9, 5.3, "Quy Trình & Thang Chấm Điểm AI", "• Quy trình: Đăng video → AI quét nội dung → Dự đoán tệp xem → Đề xuất (20% follower, 80% người mới) → Đánh giá điểm.\n\n• Thang điểm ưu tiên tăng dần:\n1. Lượt thích: 2 điểm\n2. Lượt xem hết: 4 điểm\n3. Lượt Comment: 6 điểm\n4. Lượt Chia sẻ: 8 điểm\n5. Tỷ lệ xem lại: 10 điểm (Cao nhất)")
    add_card(slide, 6.8, 1.5, 5.9, 5.3, "4 Tầng Phân Phối Video TikTok", "• Tầng 1: Đề xuất Followers + ~400 người mới.\n• Tầng 2: Mở rộng 1.000 – 5.000 người dùng mới.\n• Tầng 3: Đề xuất rộng 50.000 – 1.000.000 người dùng.\n• Tầng 4: VIRAL (>1 triệu views, tab 'For You').\n\n⚠️ Tỷ lệ đánh giá vượt chuẩn sẽ lên tầng cao hơn, dưới chuẩn sẽ bị Flop.")

    # ==================== SLIDE 28 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 28: Chiến Lược Tối Ưu Kênh TikTok Short-Form", "Module 9: Phân Phối Nội Dung Theo Nền Tảng")

    add_card(slide, 0.6, 1.5, 5.0, 5.3, "4 Bước Phân Phối Kênh", "1. Test rộng: Ra 5–10 video biến thể nội dung.\n2. Đo & đối chiếu: Lại chỉ số views, completion rate, shares, follows.\n3. Rút mẫu hiệu quả: Rút ra nhóm format/topic hợp tệp KH.\n4. Đào sâu & ổn định: Đẩy mạnh định hình profile kênh.")
    add_card(slide, 5.9, 1.5, 6.8, 5.3, "Yếu Tố Ảnh Hưởng Đề Xuất", "• Platform Signals: Watch time, Completion rate (video <5s dễ bỏ qua), Tần suất 3–5 bài/tuần, Hashtag ngách cụ thể, 30 phút đầu là 'vòng kiểm duyệt vàng'.\n\n• Creative Signals: Hook 3s đầu, Tốc độ cắt cảnh/text động, Trend & Âm thanh, CTA tự nhiên.\n\n⭐️ TikTok là Content Platform — Chất lượng nội dung là ưu tiên số 1.")

    # ==================== SLIDE 29 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 29: Thuật Toán & Phân Phối Nội Dung Trên Threads", "Module 9: Phân Phối Nội Dung Theo Nền Tảng")

    add_card(slide, 0.6, 1.5, 3.8, 5.3, "2 Khu Vực Feed Chính", "• Following: Hiển thị người đã theo dõi theo thời gian thực.\n\n• For You: Thuật toán đề xuất xu hướng, đẩy nội dung phổ biến hoặc dễ gây tranh luận.")
    add_card(slide, 4.76, 1.5, 3.8, 5.3, "Thuật Toán Threads", "• Engagement thực: Comment, Repost, Quote-post được ưu tiên > Like.\n\n• Recency: Bài viết mới được đẩy nhanh hơn.\n\n• Độ liên quan: Dựa trên sở thích, follower, dữ liệu Instagram.")
    add_card(slide, 8.93, 1.5, 3.8, 5.3, "Nội Dung 'Hot' Trên Threads", "• Câu hỏi mở, nhận định nhẹ nhàng.\n\n• Dòng ngắn gọn, vui vẻ phong cách tweet.\n\n💡 Threads chưa có ads → Cơ hội Organic Reach lớn, tạo thảo luận hiệu quả với Gen Z.")

    # ==================== SLIDE 30 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 30: Mục Tiêu Của Chiến Dịch IMC", "Module 10: Chiến Dịch IMC (Integrated Marketing Communication)")

    imc_goals = [
        ("1. Brand Awareness (Cao nhất)", "Khi ra mắt sản phẩm mới hoặc tái tung thương hiệu, cần mass awareness ngắn hạn."),
        ("2. Brand Positioning", "Chuyển đổi định vị (từ giá rẻ → cao cấp, truyền thống → hiện đại)."),
        ("3. Sales Conversion", "Hỗ trợ qua CTA mạnh, remarketing, khuyến mãi đa kênh."),
        ("4. Engagement & Loyalty", "Hỗ trợ phụ (yếu tố chính do CRM và chất lượng sản phẩm)."),
        ("5. Crisis Management", "Hỗ trợ lan tỏa; cốt lõi phụ thuộc quy trình xử lý Crisis & PR.")
    ]

    for idx, (title, desc) in enumerate(imc_goals):
        top = 1.5 + idx * 1.05
        add_card(slide, 0.6, top, 12.13, 0.95, title, desc)

    # ==================== SLIDE 31 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 31: Vị Trí Của IMC Trong Marketing Mix (6P)", "Module 10: Chiến Dịch IMC (Integrated Marketing Communication)")

    add_card(slide, 0.6, 1.5, 5.9, 5.3, "Brand Communication & Vị Trí 6P", "• Brand Communication: Quá trình truyền tải thông điệp tới target audience nhằm thay đổi nhận thức & thái độ.\n\n• Vị trí trong 6P Marketing Mix:\n(Product, Proposition, Place, Price, Package)\n→ IMC nằm ở nhánh PROMOTION / COMMUNICATION.")
    add_card(slide, 6.8, 1.5, 5.9, 5.3, "Kênh Triển Khai IMC Plan", "• ATL (Above the Line):\nTVC, Digital, Billboard, Radio... Tăng nhận diện đại chúng.\n\n• BTL (Below the Line):\nProduct display, Sales promotion, Event, Trade marketing... Thúc đẩy điểm bán.\n\n• Case Study thực tế: Chiến dịch Bia Sài Gòn '#63LÀ1' dịp Tết.")

    # ==================== SLIDE 32 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 32: Quy Trình Sáng Tạo Idea & Message", "Module 11: Sáng Tạo Ý Tưởng & Thông Điệp Truyền Thông")

    add_card(slide, 0.6, 1.4, 12.13, 1.2, "Sơ Đồ Chuyển Hóa Ý Tưởng", "Communication Task / Insight  →  Nội dung thô  →  CAMPAIGN IDEA (BIG IDEA)  →  THÔNG ĐIỆP CHIẾN DỊCH (KEY MESSAGE)")

    add_card(slide, 0.6, 2.8, 5.9, 4.0, "Campaign Idea (Big Idea)", "• Có tính bao quát, định hướng cho cả chiến dịch và từng hoạt động nhỏ.\n• Không nhất thiết phải quá hay nhưng cần RÕ RÀNG VỀ ĐỊNH HƯỚNG để team nhất quán.")
    add_card(slide, 6.8, 2.8, 5.9, 4.0, "Key Message (Thông Điệp)", "• Tiêu chí 4 ĐÚNG: Dễ hiểu — Ngắn gọn — Dễ nhớ — Thay đổi nhận thức/hành vi.\n• Là thứ thương hiệu ĐỔ TIỀN GHIM VÀO ĐẦU NGƯỜI DÙNG.")

    # ==================== SLIDE 33 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 33: Case Study Sáng Tạo: Comfort Ultra Care", "Module 11: Sáng Tạo Ý Tưởng & Thông Điệp Truyền Thông")

    add_card(slide, 0.6, 1.5, 5.9, 2.5, "1. Customer Insight", "KH không chỉ muốn mặc đồ mới đầu năm, mà tin rằng giữ áo quần luôn mới giúp giữ 'vía' may mắn ở lại lâu hơn.")
    add_card(slide, 6.8, 1.5, 5.9, 2.5, "2. Communication Task", "Comfort Ultra Care chăm sóc chuyên sâu là giải pháp kết hợp giữa thơm bền lâu & công thức Ultra Care bảo vệ áo quần tươi mới.")
    add_card(slide, 0.6, 4.2, 5.9, 2.6, "3. Big Idea (Campaign Idea)", "Lưu giữ sự tươi mới, lan tỏa năng lượng may mắn quanh năm.")
    add_card(slide, 6.8, 4.2, 5.9, 2.6, "4. Key Message (Thông điệp)", "'Áo xuân tươi màu, 'vía' may lan tỏa'")

    # ==================== SLIDE 34 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 34: Xác Định Các Giai Đoạn (Phase) Truyền Thông", "Module 11: Sáng Tạo Ý Tưởng & Thông Điệp Truyền Thông")

    add_card(slide, 0.6, 1.4, 12.13, 0.9, "Căn Cứ Phân Bổ Phase", "Dựa vào Customer Journey và Communication Objective (Communication Task) để xác định số lượng Phase (1, 2 hay 3 Phase).")

    headers = ["Kiểu chiến dịch", "Cách phân bổ Phase truyền thông"]
    rows = [
        ["Chiến dịch lớn, dài hạn, ngân sách cao", "Chia đủ 3 Phase chuẩn: Awareness → Consideration → Conversion"],
        ["Chiến dịch nhỏ, tung SP mới, nguồn lực ít", "Gộp chung 1 Phase: Awareness / Consideration / Conversion"],
        ["Dạng phổ biến, tập trung bán hàng", "Chia 2 Phase trọng tâm: Awareness → Conversion (Bỏ qua Consideration)"]
    ]
    add_table(slide, 0.6, 2.5, 12.13, 4.3, headers, rows, [4.0, 8.13])

    # ==================== SLIDE 35 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, LIGHT_BG)
    add_header(slide, "Slide 35: Tóm Tắt Toàn Bộ Mạch Kiến Thức Bài Học", "SLIDE TỔNG KẾT")

    summary_items = [
        "1. Nền tảng: 3 tầng giá trị sản phẩm & Phễu TOFU-MOFU-BOFU.",
        "2. Insight: Phân biệt Vấn đề — Pain point — Customer Insight.",
        "3. Công thức: Áp dụng AIDA, PAS, FAB & Kỹ thuật Hook tiêu đề.",
        "4. Video Content: Quy trình 6 bước & 7 dạng video chuyển đổi.",
        "5. Vận hành: Quy trình 6 bước khép kín & Chỉ số Readability/Watchability.",
        "6. Chiến lược: Mô hình 3C, Strategy Canvas & Value Matrix.",
        "7. Chuyển đổi: Tránh Mirage Content (Specific Customer + Specific Pain).",
        "8. Phân phối: Tối ưu thuật toán Facebook, TikTok & Threads.",
        "9. IMC: Định vị IMC trong Marketing Mix (6P) & Kênh ATL/BTL.",
        "10. Sáng tạo: Chuyển hóa Insight → Big Idea → Key Message & Phase."
    ]

    for idx, item in enumerate(summary_items):
        row = idx // 2
        col = idx % 2
        left = 0.6 + col * 6.2
        top = 1.4 + row * 1.1
        add_card(slide, left, top, 5.9, 0.95, "", item)

    prs.save("Giao_An_Content_Marketing.pptx")
    print(">>> Đã xuất file thành công: Giao_An_Content_Marketing.pptx")

if __name__ == "__main__":
    create_deck()
